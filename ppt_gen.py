from flask import Blueprint, request, jsonify, send_file
from pptx import Presentation
import io
import os

# 定義 Blueprint
ppt_gen_bp = Blueprint('ppt_gen', __name__)

# 定義 PowerPoint 模板檔案
template_files = {
    "A": "TempB.pptx",
    "B": "TempC.pptx"
}

# Markdown 轉 PowerPoint 的函數
def generate_pptx(md_content, template_path):
    prs = Presentation(template_path)
    slides = md_content.replace("\r", "").strip().split("\n# ")  # Remove _x000D_ artifacts
    
    for i, slide_content in enumerate(slides):
        lines = slide_content.split("\n")
        title = lines[0].replace("# ", "").strip()
        body = "\n".join(lines[1:]).strip()
        
        if i == 0:
            slide_layout = prs.slide_layouts[0]  # Title Slide Layout
        else:
            slide_layout = prs.slide_layouts[1]  # Title & Content layout
        
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if body:
            slide.shapes.placeholders[1].text = body
    
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes

# PowerPoint 生成 API
@ppt_gen_bp.route("/generate_pptx", methods=["POST"])
def generate_pptx_api():
    md_file = request.files.get("md_file")
    template_choice = request.form.get("template_file", "A")
    
    if not md_file:
        return {"error": "Missing md_file"}, 400
    
    if template_choice not in template_files:
        return {"error": "Invalid template_file choice"}, 400
    
    md_content = md_file.read().decode("utf-8")
    template_path = template_files[template_choice]
    
    if not os.path.exists(template_path):
        return {"error": "Template file not found"}, 500
    
    pptx_bytes = generate_pptx(md_content, template_path)
    return send_file(pptx_bytes, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation", as_attachment=True, download_name="output.pptx")